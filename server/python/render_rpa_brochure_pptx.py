# -*- coding: utf-8 -*-
"""Jumeau PPTX éditable de la brochure RPA — MIROIR FIDÈLE du PDF (render_rpa_brochure.py).

Reproduit chaque page aux MÊMES coordonnées que le moteur PDF (via rpa_pptx_helpers, qui
convertit l'espace ReportLab y-depuis-le-bas vers l'EMU PowerPoint). Objets natifs (texte,
formes, dégradés, scrims, images cover-crop arrondies, icônes Font Awesome en PNG net) →
le PPTX est visuellement identique au PDF. Data-driven : tout vient de `data`.

Round-trip : les boîtes de texte VERBATIM (non transformées) portent un nom « RPA::<chemin> »
relu par ingest_rpa_brochure_pptx.py. Les libellés rendus en MAJUSCULES (kickers, titres de
sous-section) ne sont pas nommés — ils s'éditent via l'éditeur de contenu structuré, pour ne
pas figer leur casse au round-trip.

E/S : lit {"data": {...}, "out": "<chemin.pptx>"} sur stdin, écrit le PPTX, renvoie {"path"}.
"""
import os
import sys
import json
import tempfile

from pptx import Presentation
from PIL import ImageFont

import rpa_pptx_helpers as H
from rpa_pptx_helpers import (E, text_line, para, rect, oval, line, grad_rect, scrim,
                              picture, picture_raw, icon, img_size, set_asset_dir, set_icon_fonts)

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
FN = os.path.join(ASSETS, "fonts")
WF = "C:/Windows/Fonts"

PW, PH = 612.0, 792.0
M = 50.0
CW = PW - 2 * M

# ── Palette (identique au PDF render_rpa_brochure.py) ──
INK = "#21303A"; INK2 = "#4A5A63"; DEEP = "#0F3B4C"; DEEP_D = "#0A2C39"; DEEP2 = "#1E6478"
GOLD = "#BF9A46"; GOLD_D = "#9A7826"; GOLD_LT = "#E8D9B0"
CREAM = "#F8F3E8"; CREAM_B = "#EBE1CB"; MIST = "#ECF2F3"; MIST_B = "#D9E5E6"
LINE = "#E1E5E2"; WHITE = "#FFFFFF"

# Codepoints Font Awesome (identiques au PDF). Clés stables du modèle de contenu.
IC = {
    "house": 0x1F3E0, "leaf": 0xF06C, "bolt": 0xF0E7, "blender": 0xF517,
    "wifi": 0xF1EB, "car": 0x1F698, "shield": 0xF3ED, "nurse": 0xF82F, "bell": 0x1F6CE,
    "key": 0x1F511, "camera": 0xF03D, "fire": 0x1F525, "steth": 0x1FA7A, "pills": 0xF484,
    "cart": 0x1F6D2, "swim": 0x1F3CA, "dumbbell": 0xF44B, "film": 0x1F39E, "book": 0x1F56E,
    "golf": 0xF450, "dice": 0x1F3B2, "church": 0xF51D, "calendar": 0xF073, "champagne": 0x1F942,
    "utensils": 0x1F374, "pin": 0xF3C5, "bus": 0x1F68D, "hospital": 0x1F3E5, "store": 0xF54E,
    "coins": 0xF51E, "phone": 0x1F57B, "envelope": 0x1F582, "check": 0xF058, "tree": 0x1F332,
    "heart": 0x1F9E1, "couch": 0xF4B8, "water": 0xF773, "masks": 0x1F3AD, "globe": 0x1F310,
    "cake": 0x1F382,
}
ICB = {"linkedin": 0xF0E1}

# Override de positions (Phase C). data["layout"][slot] = [x, y, w, h] (pt, origine bas-gauche).
LAYOUT_OV = {}


def ov(slot, x, y, w, h):
    b = LAYOUT_OV.get(slot)
    if isinstance(b, (list, tuple)) and len(b) == 4:
        return float(b[0]), float(b[1]), float(b[2]), float(b[3])
    return x, y, w, h


_mfont = {}


def _measure(ttf, size, text):
    """Largeur de texte (pt ≈ px à 72 dpi) via PIL — pour les éléments calculés (pastilles)."""
    key = (ttf, size)
    f = _mfont.get(key)
    if f is None:
        try:
            f = ImageFont.truetype(ttf, int(round(size)))
        except Exception:  # noqa: BLE001
            f = False
        _mfont[key] = f
    if not f:
        return len(text) * size * 0.5
    return f.getlength(text)


def iov(slot, cx, cy):
    """Centre d'icône depuis sa boîte override (modèle granulaire)."""
    b = LAYOUT_OV.get(slot)
    if isinstance(b, (list, tuple)) and len(b) == 4:
        return b[0] + b[2] / 2, b[1] + b[3] / 2
    return cx, cy


def ovline(slot, x1, y1, x2, y2):
    """Filet horizontal depuis sa boîte override."""
    b = LAYOUT_OV.get(slot)
    if isinstance(b, (list, tuple)) and len(b) == 4:
        return b[0], b[1], b[0] + b[2], b[1]
    return x1, y1, x2, y2


def fa(s, key, cx, cy, size, color, brand=False, slot=None):
    cp = (ICB if brand else IC).get(key)
    if cp is None:
        return None
    if slot:
        cx, cy = iov(slot, cx, cy)
    pic = icon(s, chr(cp), cx, cy, size, color, font="FAB" if brand else "FA")
    if slot and pic is not None:
        pic.name = "RPA::" + slot
    return pic


def draw_image(s, path, x, y, w, h, radius=10, dpi=240, slot=None):
    """Image cover-crop arrondie ; réserve élégante (gris MIST) si l'image est absente.
    `slot` : nomme la forme RPA::<slot> et applique l'override de position (aller-retour)."""
    if slot:
        x, y, w, h = ov(slot, x, y, w, h)
    if path and os.path.exists(path):
        shp = picture(s, path, x, y, w, h, radius=radius, dpi=dpi)
    else:
        shp = rect(s, x, y, w, h, fill=MIST, line=MIST_B, line_w=1, radius=radius)
    if slot:
        shp.name = "RPA::" + slot
    return shp


def draw_logo(s, path, x, y, w=None, h=None, anchor="bl", slot=None):
    if not (path and os.path.exists(path)):
        return 0, 0
    iw, ih = img_size(path); ar = iw / ih
    if w and not h:
        h = w / ar
    if h and not w:
        w = h * ar
    if anchor == "br":
        x -= w
    elif anchor == "tl":
        y -= h
    elif anchor == "tr":
        x -= w; y -= h
    elif anchor == "center":
        x -= w / 2; y -= h / 2
    if slot:
        x, y, w, h = ov(slot, x, y, w, h)
    pic = picture_raw(s, path, x, y, w, h)
    if slot:
        pic.name = "RPAp::" + slot
    return w, h


def kicker(s, x, y, text, color=GOLD, size=10.5, tracking=2.6, slot=None):
    if text:
        text_line(s, x, y, str(text).upper(), "Osw-SB", size, color, tracking=tracking, name=("RPAp::" + slot) if slot else None)


def title_block(s, x, y_top, kick, title_lines, title_size=30, tcolor=DEEP, rule=True, rule_w=46, kbase=None):
    kicker(s, x, y_top, kick, color=GOLD, slot=("%s.kicker" % kbase) if kbase else None)
    yy = y_top - 20
    lead = title_size * 1.02
    for i, ln in enumerate(title_lines or []):
        text_line(s, x, yy - title_size * 0.80, ln, "Osw-B", title_size, tcolor,
                  name=("RPA::%s.title.%d" % (kbase, i)) if kbase else None)
        yy -= lead
    if rule:
        if kbase:
            a, ry, b2, _ = ovline("%s.rule" % kbase, x, yy - 2, x + rule_w, yy - 2)
            line(s, a, ry, b2, ry, GOLD, 2.4).name = "RPAp::%s.rule" % kbase
        else:
            line(s, x, yy - 2, x + rule_w, yy - 2, GOLD, 2.4)
        yy -= 8
    return yy


def feature_card(s, x, y, w, h, glyph, label, desc, bg=CREAM, bdr=CREAM_B, icon_color=GOLD_D,
                 label_color=DEEP, nbase=None, slot=None):
    if not label and not desc:
        return
    if slot:
        x, y, w, h = ov(slot, x, y, w, h)  # déplacer la carte → enfants relatifs suivent
    card = rect(s, x, y, w, h, fill=bg, line=bdr, line_w=1, radius=9)
    if slot:
        card.name = "RPA::" + slot
    isl = ("%s.icon" % nbase) if nbase else None
    cxx, cyy = iov(isl, x + 26, y + h - 26) if isl else (x + 26, y + h - 26)
    oval(s, cxx, cyy, 16, fill=WHITE, line=GOLD_LT, line_w=1)
    fa(s, glyph, cxx, cyy, 15, icon_color, slot=isl)
    tx = x + 52
    text_line(s, tx, y + h - 22, label or "", "Osw-SB", 12.5, label_color,
              name=("RPA::%s.label" % nbase) if nbase else None)
    para(s, tx, y + h - 32, w - (tx - x) - 12, desc, "Sg", 9.0, INK2, 11.6,
         name=("RPA::%s.desc" % nbase) if nbase else None)


def footer(s, page_no, broker):
    y = 34.0
    line(s, M, y + 10, PW - M, y + 10, LINE, 0.8)
    bits = [broker.get("name"), broker.get("title_line"), broker.get("agency")]
    text_line(s, M, y, "  ·  ".join([b for b in bits if b]), "Sg", 7.6, INK2)
    text_line(s, PW - M, y, "%02d" % page_no, "Sg-SB", 7.6, GOLD_D, align="r")


def running_head(s, label, right_label):
    y = PH - 42
    text_line(s, M, y, str(label).upper(), "Osw-SB", 8.5, GOLD_D, tracking=2.2)
    if right_label:
        text_line(s, PW - M, y, str(right_label).upper(), "Osw-SB", 8.5, DEEP, tracking=2.2, align="r")
    line(s, M, y - 8, PW - M, y - 8, LINE, 0.8)


def _intro(s, sec, rt, num, kbase):
    running_head(s, sec.get("running", num or ""), rt)
    yb = title_block(s, M, PH - 70, sec.get("kicker", ""), sec.get("title", []), title_size=30, kbase=kbase)
    lh = 0
    if sec.get("lead"):
        para(s, M, yb - 8, CW, sec.get("lead"), "Sg", 11.5, INK, 17, name="RPA::%s.lead" % kbase)
        # hauteur approximée du paragraphe (pour le placement) — comme le PDF mesure ReportLab.
        lh = _para_h(sec.get("lead"), 11.5, 17, CW)
    return yb - 8 - lh - 16


def _para_h(text, size, leading, w):
    """Approxime la hauteur d'un paragraphe (nb de lignes × interligne)."""
    if not text:
        return 0
    avg = size * 0.50
    per = max(1, int(w / avg))
    import math
    lines = max(1, math.ceil(len(str(text)) / per))
    return lines * leading


# ════════════════════════════ PAGES ════════════════════════════
def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, PW, PH, fill=WHITE)
    return s


def page_cover(prs, d):
    A = d["assets"]; broker = d["broker"]; cov = d["content"].get("cover", {})
    s = new_slide(prs)
    hero_h = 440.0; hero_y = PH - hero_h
    draw_image(s, cov.get("hero"), 0, hero_y, PW, hero_h, radius=0, dpi=200, slot="cover.hero")
    scrim(s, 0, hero_y, PW, 180, bot_alpha=200)
    scrim(s, 0, PH - 90, PW, 90, top_alpha=150, bot_alpha=0)
    draw_logo(s, A.get("agency_logo_white"), M, PH - 34, h=30, anchor="tl", slot="cover.logo_white")
    if cov.get("pill"):
        pw = _measure(os.path.join(FN, "Oswald-600.ttf"), 9, str(cov["pill"]).upper()) + 1.4 * (len(str(cov["pill"])) - 1)
        pill_w = max(150.0, pw + 42); px = PW - M - pill_w; py = PH - 46
        pbx, pby, pbw, pbh = ov("cover.pill", px, py, pill_w, 22)
        rect(s, pbx, pby, pbw, pbh, fill=WHITE, fill_alpha=0.16, line=WHITE, line_w=1, radius=11).name = "RPAp::cover.pill"
        fa(s, "check", pbx + 15, pby + 11, 11, WHITE, slot="cover.pill_icon")
        text_line(s, pbx + 27, pby + 7, str(cov["pill"]).upper(), "Osw-SB", 9, WHITE, tracking=1.4, name="RPAp::cover.pill_text")
    if cov.get("hero_tag"):
        text_line(s, M, hero_y + 24, cov["hero_tag"], "Osw-SB", 12, WHITE, name="RPA::cover.hero_tag")
    ly = hero_y; x = M
    kicker(s, x, ly - 42, cov.get("eyebrow", ""), color=GOLD_D, size=11, slot="cover.eyebrow")
    tl = cov.get("title", [])
    yy = ly - 90
    for i, ln in enumerate(tl[:2]):
        text_line(s, x, yy, ln, "Osw-B", 46, DEEP, name="RPA::cover.title.%d" % i)
        yy -= 44
    crl = ovline("cover.rule", x, yy + 8, x + 54, yy + 8)
    line(s, crl[0], crl[1], crl[2], crl[1], GOLD, 2.6).name = "RPAp::cover.rule"
    if cov.get("subtitle"):
        para(s, x, yy - 4, CW * 0.90, cov["subtitle"], "Sg-L", 13, INK, 17, name="RPA::cover.subtitle")
    chips = cov.get("chips", [])
    sub_h = _para_h(cov.get("subtitle"), 13, 17, CW * 0.90)
    cy = max(96.0, (yy - 4 - sub_h) - 14 - 30); cx = x
    for i, chip in enumerate(chips[:3]):
        txt = chip.get("text", "")
        tw = _measure(os.path.join(WF, "seguisb.ttf"), 10, txt); cw = tw + 42
        bx, by, bw, bh = ov("cover.chips.%d" % i, cx, cy, cw, 30)
        rect(s, bx, by, bw, bh, fill=MIST, line=MIST_B, line_w=1, radius=15).name = "RPA::cover.chips.%d" % i
        fa(s, chip.get("icon"), bx + 17, by + 15, 12, GOLD_D)
        text_line(s, bx + 31, by + 10.5, txt, "Sg-SB", 10, DEEP, name="RPA::cover.chips.%d.text" % i)
        cx += cw + 11
    by = 46.0
    line(s, M, by + 34, PW - M, by + 34, LINE, 1)
    draw_logo(s, A.get("agency_logo_black"), M, by, h=26, anchor="bl", slot="cover.logo_black")
    text_line(s, PW - M, by + 16, " · ".join([b for b in [broker.get("name"), broker.get("title_line")] if b]),
              "Sg", 9.5, INK2, align="r")
    contact = "   ·   ".join([b for b in [broker.get("phone"), broker.get("email")] if b])
    text_line(s, PW - M, by + 2, contact, "Sg-SB", 9.5, DEEP, align="r")


def page_comfort(prs, d, page_no):
    sec = d["content"].get("comfort", {}); rt = d["content"].get("running_title")
    s = new_slide(prs)
    top = _intro(s, sec, rt, "01 · " + sec.get("running", ""), "comfort")
    ih = 178.0
    ib = ov("comfort.wide_image", M, top - ih, CW, ih)
    draw_image(s, sec.get("wide_image"), M, top - ih, CW, ih, radius=12, slot="comfort.wide_image")
    capn = sec.get("wide_caption") or {}
    if capn.get("text"):
        scrim(s, ib[0], ib[1], ib[2], 48, bot_alpha=150)
        fa(s, capn.get("icon"), ib[0] + 18, ib[1] + 15, 11, GOLD_LT, slot="comfort.wide_caption.icon")
        text_line(s, ib[0] + 33, ib[1] + 11, capn["text"], "Sg-SB", 9.5, WHITE, name="RPA::comfort.wide_caption.text")
    gy = top - ih - 22
    feats = sec.get("features", [])
    gap = 14.0; cardw = (CW - gap) / 2; cardh = 66.0
    for i, f in enumerate(feats[:6]):
        r = i // 2; col = i % 2
        x = M + col * (cardw + gap); y = gy - cardh - r * (cardh + gap)
        feature_card(s, x, y, cardw, cardh, f.get("icon"), f.get("label"), f.get("desc"),
                     nbase="comfort.features.%d" % i, slot="comfort.features.%d.card" % i)
    rows = (min(len(feats[:6]), 6) + 1) // 2
    note = sec.get("note") or {}
    if note.get("title"):
        ny = gy - rows * (cardh + gap) - 6
        nx, nyy, nw, nh = ov("comfort.note.card", M, ny - 46, CW, 46)
        grad_rect(s, nx, nyy, nw, nh, DEEP, DEEP2, radius=12, vertical=False).name = "RPA::comfort.note.card"
        fa(s, "check", nx + 26, nyy + nh - 23, 15, GOLD_LT, slot="comfort.note.icon")
        text_line(s, nx + 48, nyy + nh - 19, note["title"], "Osw-SB", 12.5, WHITE, name="RPA::comfort.note.title")
        if note.get("sub"):
            text_line(s, nx + 48, nyy + nh - 32, note["sub"], "Sg", 9.5, GOLD_LT, name="RPA::comfort.note.sub")
    footer(s, page_no, d["broker"])


def page_security(prs, d, page_no):
    sec = d["content"].get("security", {}); rt = d["content"].get("running_title")
    s = new_slide(prs)
    top = _intro(s, sec, rt, "02 · " + sec.get("running", ""), "security")
    panel_w = CW * 0.605; img_w = CW - panel_w - 16; panel_h = 236.0
    px, py, pw, ph = ov("security.panel", M, top - panel_h, panel_w, panel_h)
    grad_rect(s, px, py, pw, ph, DEEP, DEEP_D, radius=14, vertical=True).name = "RPA::security.panel"
    if sec.get("panel_title"):
        text_line(s, px + 22, py + ph - 26, str(sec["panel_title"]).upper(), "Osw-SB", 12, GOLD_LT, tracking=1.6, name="RPAp::security.panel_title")
    iy = py + ph - 50
    for i, it in enumerate(sec.get("panel_items", [])[:5]):
        fa(s, it.get("icon"), px + 30, iy - 2, 14, GOLD, slot="security.panel_items.%d.icon" % i)
        para(s, px + 50, iy + 6, pw - 50 - 18, it.get("text"), "Sg", 10, WHITE, 13.0,
             name="RPA::security.panel_items.%d.text" % i)
        iy -= 37
    draw_image(s, sec.get("panel_image"), M + panel_w + 16, top - panel_h, img_w, panel_h, radius=14, slot="security.panel_image")
    if sec.get("panel_caption"):
        scrim(s, M + panel_w + 16, top - panel_h, img_w, 70, bot_alpha=150)
        text_line(s, M + panel_w + 16 + 12, top - panel_h + 12, sec["panel_caption"], "Sg-SB", 9, WHITE,
                  name="RPA::security.panel_caption")
    sy = top - panel_h - 26
    kicker(s, M, sy, sec.get("services_kicker", ""), color=GOLD_D, size=10.5, slot="security.services_kicker")
    if sec.get("services_title"):
        text_line(s, M, sy - 26, str(sec["services_title"]).upper(), "Osw-B", 21, DEEP, name="RPAp::security.services_title")
    svcs = sec.get("services", [])
    gap = 14.0; cw = (CW - 2 * gap) / 3; chh = 92.0; sc_y = sy - 40
    for i, sv in enumerate(svcs[:3]):
        x = M + i * (cw + gap); y = sc_y - chh
        bx, by, bw, bh = ov("security.services.%d.card" % i, x, y, cw, chh)
        rect(s, bx, by, bw, bh, fill=CREAM, line=CREAM_B, line_w=1, radius=10).name = "RPA::security.services.%d.card" % i
        scx, scy = iov("security.services.%d.icon" % i, bx + 26, by + bh - 26)
        oval(s, scx, scy, 17, fill=WHITE, line=GOLD_LT, line_w=1)
        fa(s, sv.get("icon"), scx, scy, 16, GOLD_D, slot="security.services.%d.icon" % i)
        text_line(s, bx + 52, by + bh - 30, sv.get("label", ""), "Osw-SB", 13, DEEP,
                  name="RPA::security.services.%d.label" % i)
        para(s, bx + 16, by + bh - 48, bw - 30, sv.get("desc"), "Sg", 9.2, INK2, 11.8,
             name="RPA::security.services.%d.desc" % i)
    footer(s, page_no, d["broker"])


def page_amenities(prs, d, page_no):
    sec = d["content"].get("amenities", {}); rt = d["content"].get("running_title")
    s = new_slide(prs)
    top = _intro(s, sec, rt, "03 · " + sec.get("running", ""), "amenities")
    gallery = sec.get("gallery", [])

    def cap(x, y, w, item, idx):
        if not item.get("caption"):
            return
        scrim(s, x, y, w, 42, bot_alpha=165)
        fa(s, item.get("icon"), x + 15, y + 13, 11, GOLD_LT, slot="amenities.gallery.%d.icon" % idx)
        text_line(s, x + 30, y + 9, item["caption"], "Sg-SB", 9.5, WHITE, name="RPA::amenities.gallery.%d.caption" % idx)
    gap = 12.0; big_w = CW * 0.60; big_h = 196.0; sm_w = CW - big_w - gap; sm_h = (196.0 - gap) / 2
    g = (gallery + [{}] * 6)[:6]
    bx, by = M, top - big_h
    b = ov("amenities.gallery.0.image", bx, by, big_w, big_h)
    draw_image(s, g[0].get("image"), bx, by, big_w, big_h, radius=12, slot="amenities.gallery.0.image"); cap(b[0], b[1], b[2], g[0], 0)
    rx = M + big_w + gap
    b = ov("amenities.gallery.1.image", rx, top - sm_h, sm_w, sm_h)
    draw_image(s, g[1].get("image"), rx, top - sm_h, sm_w, sm_h, radius=12, slot="amenities.gallery.1.image"); cap(b[0], b[1], b[2], g[1], 1)
    b = ov("amenities.gallery.2.image", rx, top - big_h, sm_w, sm_h)
    draw_image(s, g[2].get("image"), rx, top - big_h, sm_w, sm_h, radius=12, slot="amenities.gallery.2.image"); cap(b[0], b[1], b[2], g[2], 2)
    r2y = by - gap; cw3 = (CW - 2 * gap) / 3; h3 = 150.0
    for i in range(3):
        item = g[3 + i]
        x = M + i * (cw3 + gap); y = r2y - h3
        sl = "amenities.gallery.%d.image" % (3 + i); b = ov(sl, x, y, cw3, h3)
        draw_image(s, item.get("image"), x, y, cw3, h3, radius=12, slot=sl); cap(b[0], b[1], b[2], item, 3 + i)
    pillars = sec.get("pillars", [])
    sy = r2y - h3 - 22; gap = 14.0; cw = (CW - 2 * gap) / 3; chh = 66.0
    for i, p in enumerate(pillars[:3]):
        x = M + i * (cw + gap); y = sy - chh
        feature_card(s, x, y, cw, chh, p.get("icon"), p.get("label"), p.get("desc"), bg=MIST, bdr=MIST_B,
                     nbase="amenities.pillars.%d" % i, slot="amenities.pillars.%d.card" % i)
    footer(s, page_no, d["broker"])


def page_life(prs, d, page_no):
    sec = d["content"].get("life", {}); rt = d["content"].get("running_title")
    s = new_slide(prs)
    top = _intro(s, sec, rt, "04 · " + sec.get("running", ""), "life")
    gap = 14.0; cw = (CW - 2 * gap) / 3; ch_img = 82.0; chh = 150.0
    for i, ev in enumerate(sec.get("events", [])[:3]):
        x = M + i * (cw + gap); y = top - chh
        bx, by, bw, bh = ov("life.events.%d.card" % i, x, y, cw, chh)
        rect(s, bx, by, bw, bh, fill=WHITE, line=CREAM_B, line_w=1, radius=10).name = "RPA::life.events.%d.card" % i
        draw_image(s, ev.get("image"), bx, by + bh - ch_img, bw, ch_img, radius=10)
        fa(s, ev.get("icon"), bx + 18, by + bh - ch_img - 14, 14, GOLD_D, slot="life.events.%d.icon" % i)
        text_line(s, bx + 34, by + bh - ch_img - 18, ev.get("label", ""), "Osw-SB", 12.5, DEEP,
                  name="RPA::life.events.%d.label" % i)
        para(s, bx + 14, by + bh - ch_img - 30, bw - 26, ev.get("desc"), "Sg", 9.0, INK2, 11.6,
             name="RPA::life.events.%d.desc" % i)
    qy = top - chh - 24
    kicker(s, M, qy, sec.get("neighborhood_kicker", ""), color=GOLD_D, size=10.5, slot="life.neighborhood_kicker")
    if sec.get("neighborhood_title"):
        text_line(s, M, qy - 26, str(sec["neighborhood_title"]).upper(), "Osw-B", 21, DEEP, name="RPAp::life.neighborhood_title")
    qcards = sec.get("neighborhood", [])
    gap = 14.0; cw2 = (CW - gap) / 2; chh2 = 58.0; gy = qy - 38
    for i, q in enumerate(qcards[:4]):
        r = i // 2; col = i % 2
        x = M + col * (cw2 + gap); y = gy - chh2 - r * (chh2 + gap)
        feature_card(s, x, y, cw2, chh2, q.get("icon"), q.get("label"), q.get("desc"), bg=MIST, bdr=MIST_B,
                     nbase="life.neighborhood.%d" % i, slot="life.neighborhood.%d.card" % i)
    fin = sec.get("finance") or {}
    if fin.get("title"):
        nrows = (min(len(qcards), 4) + 1) // 2
        fy = gy - nrows * (chh2 + gap) - 8; fh = 58.0
        fx, fyy, fw, fhh = ov("life.finance.card", M, fy - fh, CW, fh)
        grad_rect(s, fx, fyy, fw, fhh, GOLD_D, GOLD, radius=14, vertical=False).name = "RPA::life.finance.card"
        ficx, ficy = iov("life.finance.icon", fx + 34, fyy + fhh / 2)
        oval(s, ficx, ficy, 19, fill=WHITE, fill_alpha=0.92)
        fa(s, fin.get("icon", "coins"), ficx, ficy, 18, GOLD_D, slot="life.finance.icon")
        text_line(s, fx + 66, fyy + fhh - 24, str(fin["title"]).upper(), "Osw-B", 16, WHITE, name="RPAp::life.finance.title")
        para(s, fx + 66, fyy + fhh - 30, fw - 150, fin.get("text"), "Sg-SB", 9.6, WHITE, 12, name="RPA::life.finance.text")
    footer(s, page_no, d["broker"])


def page_contact(prs, d, page_no):
    A = d["assets"]; broker = d["broker"]; sec = d["content"].get("contact", {})
    s = new_slide(prs)
    txt_ref = PH - 250; band_bottom = PH - 308; band_h = PH - band_bottom
    draw_image(s, sec.get("hero"), 0, band_bottom, PW, band_h, radius=0, dpi=200, slot="contact.hero")
    scrim(s, 0, band_bottom, PW, band_h, top_alpha=120, bot_alpha=225)
    draw_logo(s, A.get("agency_logo_white"), M, PH - 34, h=28, anchor="tl", slot="contact.logo_white")
    kicker(s, M, txt_ref + 150, sec.get("kicker", ""), color=GOLD_LT, size=11, slot="contact.kicker")
    tl = sec.get("title", [])
    yy = txt_ref + 108
    for i, ln in enumerate(tl[:2]):
        text_line(s, M, yy, ln, "Osw-B", 38, WHITE, name="RPA::contact.title.%d" % i)
        yy -= 38
    krl = ovline("contact.rule", M, yy + 6, M + 54, yy + 6)
    line(s, krl[0], krl[1], krl[2], krl[1], GOLD, 2.6).name = "RPAp::contact.rule"
    if sec.get("cta"):
        para(s, M, yy - 6, CW * 0.74, sec.get("cta"), "Sg-L", 13, WHITE, 18, name="RPA::contact.cta")
    cardx = M; cardw = CW * 0.57; cardy = 148.0; cardh = band_bottom - 148 - 26
    cardx, cardy, cardw, cardh = ov("contact.card", cardx, cardy, cardw, cardh)
    rect(s, cardx, cardy, cardw, cardh, fill=WHITE, line=LINE, line_w=1.2, radius=14).name = "RPA::contact.card"
    gsx, gsy, gsw, gsh = ov("contact.gold_strip", cardx, cardy + cardh - 6, cardw, 6)
    rect(s, gsx, gsy, gsw, gsh, fill=GOLD, radius=3).name = "RPAp::contact.gold_strip"
    ix = cardx + 28; iy = cardy + cardh - 42
    text_line(s, ix, iy, str(broker.get("name", "")).upper(), "Osw-B", 26, DEEP, name="RPAp::contact.name")
    text_line(s, ix, iy - 17, broker.get("title_line", ""), "Sg-SB", 10.5, GOLD_D, name="RPAp::contact.designation")
    agency = "  ·  ".join([b for b in [broker.get("agency"), broker.get("company")] if b])
    text_line(s, ix, iy - 32, agency, "Sg", 9.5, INK2, name="RPAp::contact.agency")
    line(s, ix, iy - 46, cardx + cardw - 28, iy - 46, LINE, 1)
    rows = []
    if broker.get("phone"):
        rows.append(("phone", broker["phone"], False))
    if broker.get("email"):
        rows.append(("envelope", broker["email"], False))
    if broker.get("linkedin"):
        rows.append(("linkedin", broker["linkedin"], True))
    ry = iy - 68
    for ri, (g, txt, brand) in enumerate(rows):
        rcx, rcy = iov("contact.row.%d.icon" % ri, ix + 13, ry)
        oval(s, rcx, rcy, 13, fill=DEEP)
        fa(s, g, rcx, rcy, 12, WHITE, brand=brand, slot="contact.row.%d.icon" % ri)
        text_line(s, ix + 36, ry - 4.5, txt, "Sg-SB", 11.5, INK, name="RPAp::contact.row.%d.text" % ri)
        ry -= 31
    EXP_W = 118.0
    draw_logo(s, A.get("agency_logo_black"), ix, cardy + 20, w=EXP_W, anchor="bl", slot="contact.card_logo")
    qr = A.get("qr")
    if qr and os.path.exists(qr):
        qrs = 70.0; qx = cardx + cardw - 28 - qrs; qy = cardy + 18
        qbx, qby, qbw, qbh = ov("contact.qr", qx, qy, qrs, qrs)
        picture_raw(s, qr, qbx, qby, qbw, qbh).name = "RPAp::contact.qr"
        if broker.get("linkedin_label"):
            text_line(s, qbx + qbw / 2, qby - 9, broker["linkedin_label"], "Sg", 7.4, INK2, align="c", name="RPAp::contact.qr_label")
    if sec.get("disclaimer"):
        para(s, M, 92, CW * 0.50, sec["disclaimer"], "Sg", 7.4, INK2, 9.6, name="RPA::contact.disclaimer")
    footer(s, page_no, broker)
    sp = A.get("broker_hero")
    if sp and os.path.exists(sp):
        iw, ih = img_size(sp); ar = iw / ih; sp_h = 312.0; sp_w = sp_h * ar
        sp_x = PW - sp_w + 26; sp_y = 0
        spx, spy, spw, sph = ov("contact.broker_hero", sp_x, sp_y, sp_w, sp_h)
        picture_raw(s, sp, spx, spy, spw, sph).name = "RPAp::contact.broker_hero"
        if A.get("company_logo") and os.path.exists(A["company_logo"]):
            pv_w = EXP_W * 1.5
            draw_logo(s, A["company_logo"], spx + spw * 0.56, spy + sph * 0.245, w=pv_w, anchor="center", slot="contact.company_logo")


def render(data, out):
    broker = data.setdefault("broker", {})
    broker.setdefault("title_line", " ".join([s for s in [broker.get("title"), broker.get("subtitle")] if s]))
    data.setdefault("assets", {})
    data.setdefault("content", {})
    global LAYOUT_OV
    LAYOUT_OV = data.get("layout") or {}
    # Modèle GRANULAIRE : chaque TEXTE nommé est placé à sa propre boîte capturée (text_line/para
    # consultent POS par nom). Les slots non-texte (.card/.image) n'y sont jamais cherchés → sans effet.
    H.set_pos(LAYOUT_OV)

    cache = os.path.join(tempfile.gettempdir(), "softimmo_rpa_pptx")
    set_asset_dir(cache)
    set_icon_fonts(os.path.join(FN, "fa-solid-900.ttf"), os.path.join(FN, "fa-brands-400.ttf"))

    prs = Presentation()
    prs.slide_width = E(PW)
    prs.slide_height = E(PH)
    page_cover(prs, data)
    page_comfort(prs, data, 2)
    page_security(prs, data, 3)
    page_amenities(prs, data, 4)
    page_life(prs, data, 5)
    page_contact(prs, data, 6)
    try:
        cov = data["content"].get("cover", {})
        prs.core_properties.title = (cov.get("title") or ["Brochure RPA"])[0]
        prs.core_properties.author = broker.get("name", "Softimmo")
    except Exception:  # noqa: BLE001
        pass
    prs.save(out)
    return out


def main():
    try:
        raw = sys.stdin.buffer.read().decode("utf-8") or "{}"
        payload = json.loads(raw)
        data = payload.get("data") or {}
        out = payload.get("out")
        if not out:
            print(json.dumps({"error": "out requis"})); return
        render(data, out)
        print(json.dumps({"path": out}))
    except Exception as e:  # noqa: BLE001
        import traceback
        print(json.dumps({"error": str(e), "trace": traceback.format_exc()[-900:]}))


if __name__ == "__main__":
    main()
