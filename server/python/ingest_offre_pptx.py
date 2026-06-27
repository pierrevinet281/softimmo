# -*- coding: utf-8 -*-
"""Ingestion d'un PPTX d'offre édité → contenu de l'offre (aller-retour, docs/09).

Relit les formes NOMMÉES (« OFF::<clé>::<type>::<partie> ») écrites par render_offre_pptx.py
et reconstruit le contenu prêt au rendu : ordre des sections (= ordre des diapos), titre +
corps de chaque section selon son type (groupes/liste/étapes/témoignages/corps), et images
des sections « asset » (extraites et enregistrées). Déterministe, sans IA.

E/S : lit {"pptx": "<chemin>", "images_dir": "<dossier>"} sur stdin ;
renvoie {"content": {doc_title, subtitle, sections:[...], <clé>:{...}, ...}}.
"""
import os
import sys
import json

from pptx import Presentation

CANON = ["why", "guarantee", "services", "marketing", "opportunities",
         "timeline", "fees", "value_add", "testimonials", "next_steps"]


def _paras(shape):
    """Liste de (niveau, texte) des paragraphes non vides d'une forme texte."""
    out = []
    try:
        if not shape.has_text_frame:
            return out
        for p in shape.text_frame.paragraphs:
            txt = "".join(r.text for r in p.runs).strip()
            if txt:
                out.append((p.level or 0, txt))
    except Exception:  # noqa: BLE001
        pass
    return out


def _text(shape):
    return "\n".join(t for _, t in _paras(shape)).strip()


def _save_pic(shape, key, images_dir):
    try:
        img = shape.image
    except Exception:  # noqa: BLE001
        return None
    if not img.blob:
        return None
    os.makedirs(images_dir, exist_ok=True)
    ext = (img.ext or "png").lstrip(".")
    dest = os.path.join(images_dir, "offre_%s.%s" % (key, ext))
    with open(dest, "wb") as f:
        f.write(img.blob)
    return dest


def _split2(s, sep=" — "):
    if sep in s:
        a, b = s.split(sep, 1)
        return a.strip(), b.strip()
    return s.strip(), ""


def _build_section(stype, head, body_paras):
    """Reconstruit l'objet section selon son type structurel."""
    if stype in ("groups", "subgroups"):
        groups = []
        for lvl, txt in body_paras:
            if lvl <= 0:
                groups.append({"label": txt, "items": []})
            else:
                if not groups:
                    groups.append({"label": "", "items": []})
                groups[-1]["items"].append(txt)
        field = "subgroups" if stype == "subgroups" else "groups"
        return {"heading": head, field: groups, "kind": "groups"}
    if stype == "list":
        return {"heading": head, "items": [t for _, t in body_paras], "kind": "list"}
    if stype == "steps":
        steps = []
        for _, t in body_paras:
            lab, txt = _split2(t)
            steps.append({"label": lab, "text": txt})
        return {"heading": head, "steps": steps}
    if stype == "testi":
        items = []
        for _, t in body_paras:
            q = t
            a = ""
            if t.startswith("«"):
                inner = t.lstrip("«").strip()
                if "»" in inner:
                    q, rest = inner.split("»", 1)
                    q = q.strip()
                    a = rest.replace("—", "", 1).strip() if "—" in rest else rest.strip()
            else:
                q, a = _split2(t)
            items.append({"quote": q, "author": a})
        return {"heading": head, "items": items}
    # body
    texts = [t for _, t in body_paras]
    sec = {"heading": head, "body": texts[0] if texts else ""}
    if len(texts) > 1:
        sec["note"] = texts[-1]
        if len(texts) > 2:
            sec["body"] = "\n".join(texts[:-1])
    return sec


def build_content(prs, images_dir=None):
    content = {}
    sections = []
    seen = set()
    for slide in prs.slides:
        # regrouper les formes par clé de section présentes sur la diapo
        heads, bodies, stypes, asset_imgs, caps = {}, {}, {}, {}, {}
        for sh in slide.shapes:
            name = getattr(sh, "name", "") or ""
            if name == "OFF::__title":
                ps = _paras(sh)
                if ps:
                    content["doc_title"] = ps[0][1]
                    if len(ps) >= 2:
                        content["subtitle"] = ps[1][1]
                continue
            if not name.startswith("OFF::") or name.startswith("OFF::__"):
                continue
            parts = name.split("::")
            if len(parts) < 4:
                continue
            _, key, stype, part = parts[0], parts[1], parts[2], parts[3]
            stypes[key] = stype
            if part == "head":
                heads[key] = _text(sh)
            elif part == "body":
                bodies[key] = _paras(sh)
            elif part == "img":
                if images_dir:
                    asset_imgs[key] = _save_pic(sh, key, images_dir)
            elif part == "cap":
                caps[key] = _text(sh)
        for key in list(stypes.keys()):
            if key in seen:
                continue
            seen.add(key)
            stype = stypes[key]
            if stype == "asset":
                content[key] = {"kind": "asset", "image": asset_imgs.get(key), "caption": caps.get(key, "")}
                sections.append({"key": key, "custom": True, "kind": "asset", "hidden": False})
            else:
                content[key] = _build_section(stype, heads.get(key, ""), bodies.get(key, []))
                is_custom = key not in CANON
                sections.append({"key": key, "hidden": False, "custom": is_custom,
                                 "kind": content[key].get("kind")})
    content["sections"] = sections
    return content


def _emit(obj):
    sys.stdout.buffer.write(json.dumps(obj, ensure_ascii=False).encode("utf-8"))
    sys.stdout.buffer.flush()


def main():
    try:
        payload = json.loads(sys.stdin.buffer.read().decode("utf-8") or "{}")
        pptx = payload.get("pptx")
        if not pptx:
            _emit({"error": "pptx requis"}); return
        prs = Presentation(pptx)
        content = build_content(prs, images_dir=payload.get("images_dir"))
        _emit({"content": content})
    except Exception as e:  # noqa: BLE001
        import traceback
        _emit({"error": str(e), "trace": traceback.format_exc()[-700:]})


if __name__ == "__main__":
    main()
