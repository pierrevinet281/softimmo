# -*- coding: utf-8 -*-
"""Jumeau PowerPoint éditable de la brochure (docs/09).

Produit un .pptx FIDÈLE et MODIFIABLE, aux MÊMES coordonnées que le PDF (espace 540×720 pt
= 7,5×10 po, natif PowerPoint). Le courtier ajuste les positions dans PowerPoint, puis
`extract_pptx_layout.py` permet de reporter les changements dans `render_brochure.py`.

Déterministe, sans IA (CLAUDE.md §3). E/S : lit {"data": {...}, "out": "<chemin.pptx>"} sur
stdin, écrit le PPTX et renvoie {"path"}.
"""
import sys
import os
import io
import json

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from brochure_layout import load_layout
from brochure_slots import (
    ovr, title_lines, address_lines, broker_lines, broker2_lines, grid_cell, grid_slot,
)

# Mise en page courante (positions PPTX), peuplée par render() — voir brochure_layout.py.
# Régions par défaut + tout override de POSITION par slot capté à l'aller-retour. CHAQUE forme est
# nommée « STD::<slot> » (texte éditable + position) ou « STDp::<slot> » (position seule), miroir
# exact du moteur PDF (render_brochure.py) — voir brochure_slots.py.
LAYOUT = {}

try:
    from reportlab.graphics.barcode import qr
    from PIL import Image
except Exception:  # noqa: BLE001
    qr = None
    Image = None

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")


def asset(*p):
    return os.path.join(ASSETS, *p)


FONT = "Segoe UI"

# Thèmes (mêmes couleurs que le moteur PDF render_brochure.py).
THEMES = {
    "unifamilial": {
        "banner": "314897", "title_fg": "FFFFFF", "sub_fg": "FFFFFF", "title_upper": False, "luxe": False,
        "label_bg": "314897", "label_fg": "FFFFFF", "value_bg": "D7DEEE", "value_fg": "1A1A1A",
        "rule": "314897", "price_bg": "E2231A", "price_fg": "FFFFFF", "bar": "E2231A", "qr": "E2231A",
        "p2_banner": "E2231A", "p2_title_fg": "FFFFFF", "desc_bg": "C6D1E6",
        "th_bg": "314897", "th_fg": "FFFFFF", "row": "EEF1F7", "row_alt": "D7DEEE", "row_fg": "1A1A1A",
        "logo": asset("unifamilial", "exp_logo_white.png"), "medal": asset("unifamilial", "certificat.png"),
        "hero": asset("unifamilial", "superpierre.png"),
    },
    "luxe": {
        "banner": "221F1C", "title_fg": "B79A5B", "sub_fg": "FFFFFF", "title_upper": True, "luxe": True,
        "label_bg": "B79A5B", "label_fg": "FFFFFF", "value_bg": "F2EDE2", "value_fg": "1A1A1A",
        "rule": "B79A5B", "price_bg": "221F1C", "price_fg": "FFFFFF", "bar": "B79A5B", "qr": "9C8246",
        "p2_banner": "9C8246", "p2_title_fg": "FFFFFF", "desc_bg": "F2EDE2",
        "th_bg": "9C8246", "th_fg": "FFFFFF", "row": "F7F3EA", "row_alt": "F2EDE2", "row_fg": "1A1A1A",
        "logo": asset("luxe", "exp_luxury_white.png"), "medal": None,
        "hero": asset("luxe", "superpierre_luxury.png"),
    },
}


def _theme_pptx(banner, price, desc, value, p2=None):
    """Thème « bannière simple » (logo + titre, sans médaille) — RPA/commercial/industriel."""
    return {
        "banner": banner, "title_fg": "FFFFFF", "sub_fg": "FFFFFF", "title_upper": False, "luxe": False,
        "label_bg": banner, "label_fg": "FFFFFF", "value_bg": value, "value_fg": "1A1A1A",
        "rule": banner, "price_bg": price, "price_fg": "FFFFFF", "bar": price, "qr": price,
        "p2_banner": p2 or banner, "p2_title_fg": "FFFFFF", "desc_bg": desc,
        "th_bg": banner, "th_fg": "FFFFFF", "row": "F1F4F5", "row_alt": value, "row_fg": "1A1A1A",
        "logo": asset("unifamilial", "exp_logo_white.png"), "medal": None,
        "hero": asset("unifamilial", "superpierre.png"),
    }


THEMES["rpa"] = _theme_pptx("2E6E5E", "C25E3A", "DCE9E5", "DCE9E5")
THEMES["commercial"] = _theme_pptx("243B53", "C0392B", "DCE2EC", "DCE2EC")
THEMES["industriel"] = _theme_pptx("37474F", "E07B2C", "E0E4E6", "E0E4E6")

PLACEHOLDER = "E9EDF3"


def _rgb(h):
    return RGBColor.from_string(h)


def _name(sp, name):
    """Nomme la forme (nom stable du gabarit → round-trip self-service via NAME_MAP)."""
    if sp is not None and name:
        try:
            sp.name = name
        except Exception:  # noqa: BLE001
            pass
    return sp


def _rect(slide, x, y, w, h, fill=None, line=None, text=None, size=12, bold=False,
          color="000000", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, padx=2.0, name=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x), Pt(y), Pt(w), Pt(h))
    try:
        sp.shadow.inherit = False
    except Exception:  # noqa: BLE001
        pass
    if fill:
        sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = _rgb(line); sp.line.width = Pt(0.75)
    else:
        sp.line.fill.background()
    if text is not None:
        tf = sp.text_frame; tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = Pt(1); tf.margin_left = tf.margin_right = Pt(padx)
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT; r.font.color.rgb = _rgb(color)
    return _name(sp, name)


def _textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, wrap=True, name=None):
    """lines : liste de (texte, taille, gras, couleur_hex, alignement)."""
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.margin_top = tf.margin_bottom = Pt(0); tf.margin_left = tf.margin_right = Pt(0)
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align or PP_ALIGN.LEFT
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bool(bold); r.font.name = FONT
        if color:
            r.font.color.rgb = _rgb(color)
    return _name(tb, name)


PLACEHOLDER_IMG = asset("_placeholder.png")  # réserve grise nommée (remplaçable dans PowerPoint)


def _pic(slide, path, x, y, w, h, name=None):
    """Image au cadre exact. Emplacement vide = image de réserve NOMMÉE → le courtier peut
    « Changer l'image » dans PowerPoint (le nom de forme est conservé → sync de l'image)."""
    src = path if (path and os.path.exists(path)) else PLACEHOLDER_IMG
    if os.path.exists(src):
        _name(slide.shapes.add_picture(src, Pt(x), Pt(y), Pt(w), Pt(h)), name)
    else:
        _rect(slide, x, y, w, h, fill=PLACEHOLDER, text="image", size=8, color="5A5A5A", name=name)


def _pic_fit(slide, path, x, y, w=None, h=None, name=None):
    """Image en préservant l'aspect (largeur OU hauteur imposée), ancrée en (x, y)."""
    if path and os.path.exists(path):
        _name(slide.shapes.add_picture(path, Pt(x), Pt(y),
                                       width=Pt(w) if w else None, height=Pt(h) if h else None), name)


def _qr_stream(url, dark="E2231A"):
    """Génère un PNG de QR code (matrice ReportLab rasterisée via PIL) → BytesIO."""
    if not (url and qr and Image):
        return None
    w = qr.QrCodeWidget(str(url)); w.draw()
    m = w.qr.modules; n = len(m); scale, border = 10, 4
    side = (n + 2 * border) * scale
    rgb = tuple(int(dark[i:i + 2], 16) for i in (0, 2, 4))
    im = Image.new("RGB", (side, side), "white"); px = im.load()
    for r in range(n):
        for cI in range(n):
            if m[r][cI]:
                for dy in range(scale):
                    for dx in range(scale):
                        px[(cI + border) * scale + dx, (r + border) * scale + dy] = rgb
    buf = io.BytesIO(); im.save(buf, "PNG"); buf.seek(0)
    return buf


# ── Dessin par SLOT (override de position + nommage aller-retour — miroir du moteur PDF) ──
def _pos(slot, x, y, w, h):
    return ovr(LAYOUT, slot, x, y, w, h)


def _line(slide, slot, box, text, size, bold, color, align=PP_ALIGN.LEFT, editable=True,
          anchor=MSO_ANCHOR.MIDDLE, padx=0.0):
    """Ligne de texte centrée dans la boîte ; vide → non dessinée (miroir du PDF `_line`).
    `editable=True` → « STD:: » (texte relu) ; sinon « STDp:: » (position seule)."""
    if text is None or text == "":
        return None
    x, y, w, h = _pos(slot, *box)
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Pt(0); tf.margin_left = tf.margin_right = Pt(padx)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bool(bold); r.font.name = FONT; r.font.color.rgb = _rgb(color)
    return _name(tb, ("STD::" if editable else "STDp::") + slot)


def _cell(slide, slot, box, text, fill, fg, size=12):
    """Cellule de grille (rectangle plein + texte gauche) — position seule (STDp::)."""
    x, y, w, h = _pos(slot, *box)
    return _rect(slide, x, y, w, h, fill=fill, text=text, size=size, color=fg,
                 align=PP_ALIGN.LEFT, padx=8, name="STDp::" + slot)


def _broker_url(d, broker):
    u = d.get("listing_url") or broker.get("web") or broker.get("website")
    if u and not str(u).startswith(("http://", "https://")):
        u = "https://" + str(u)
    return u


# ───────────────────────────── Page 1 ─────────────────────────────
def slide1(prs, d, th):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    img = d.get("images", {}) or {}; broker = d.get("broker", {})
    title = d.get("title", "")
    if th["title_upper"]:
        title = title.upper()

    L = LAYOUT
    _rect(s, *_pos("banner", *L["banner"]), fill=th["banner"], name="STDp::banner")  # bannière
    if th["luxe"]:
        tl = title_lines(L, key="luxe_title")  # titre or + ville + résumé (lignes individuelles)
        _line(s, "title", tl["title"], title, 20, True, th["title_fg"])
        _line(s, "city", tl["city"], d.get("city", ""), 12, False, th["sub_fg"])
        _line(s, "summary_line", tl["summary_line"], d.get("summary_line", ""), 12, False, th["sub_fg"])
        lkx, lky, lkw, lkh = _pos("luxe_lock", *L["luxe_lock"])
        _pic_fit(s, img.get("logo") or th.get("logo"), lkx, lky, w=lkw, name="STDp::luxe_lock")
    else:
        lx, ly, lw, lh = _pos("logo", *L["logo"])
        _pic_fit(s, img.get("logo") or th.get("logo"), lx, ly, h=lh, name="STDp::logo")
        tl = title_lines(L)  # titre / ville / résumé (lignes individuelles, déplaçables)
        _line(s, "title", tl["title"], title, 20, True, th["title_fg"])
        _line(s, "city", tl["city"], d.get("city", ""), 12, False, th["sub_fg"])
        _line(s, "summary_line", tl["summary_line"], d.get("summary_line", ""), 12, False, th["sub_fg"])

    _pic(s, img.get("hero"), *_pos("hero", *L["hero"]), name="STDp::hero")   # photo
    _pic(s, img.get("map"), *_pos("map", *L["map"]), name="STDp::map")       # carte
    if th.get("medal"):  # médaille (déborde le haut de la bannière)
        mx, my, mw, mh = _pos("medal", *L["medal"])
        _pic_fit(s, img.get("medal") or th["medal"], mx, my, w=mw, name="STDp::medal")

    al = address_lines(L)  # adresse + MLS (lignes individuelles)
    _line(s, "address", al["address"], d.get("address", ""), 16, True, "1A1A1A")
    if d.get("mls"):
        _line(s, "mls", al["mls"], "MLS : %s" % d["mls"], 11, False, "5A5A5A")
    rlx, rly, rlw, rlh = _pos("rule", *L["rule"])
    _rect(s, rlx, rly, rlw, max(rlh, 1.6), fill=th["rule"], name="STDp::rule")  # filet

    # Grille de spécifications — CHAQUE cellule est une forme nommée (round-trip granulaire).
    left = d.get("specs_left", []); right = d.get("specs_right", [])
    for i in range(5):
        for ci, specs in enumerate((left, right)):
            if i < len(specs):
                lab = specs[i][0]; val = "" if specs[i][1] is None else str(specs[i][1])
                _cell(s, grid_slot(i, ci, "label"), grid_cell(L, i, ci, "label"), lab, th["label_bg"], th["label_fg"])
                _cell(s, grid_slot(i, ci, "value"), grid_cell(L, i, ci, "value"), val, th["value_bg"], th["value_fg"])

    price = d.get("price")  # bloc prix (texte éditable)
    ptxt = d.get("price_text") or (("Prix : %s $" % format(int(price), ",d").replace(",", " ")) if price else "Prix sur demande")
    _rect(s, *_pos("price_text", *L["price"]), fill=th["price_bg"], text=ptxt, size=28, bold=True,
          color=th["price_fg"], align=PP_ALIGN.CENTER, name="STD::price_text")

    _pic(s, broker.get("photo") or asset("broker", "portrait.png"), *_pos("broker_photo", *L["broker_photo"]),
         name="STDp::broker_photo")
    bl = broker_lines(L)  # 5 lignes (texte issu du profil → position seule)
    _line(s, "broker.name", bl["broker.name"], broker.get("name", ""), 15, True, "1A1A1A", editable=False)
    _line(s, "broker.title", bl["broker.title"], broker.get("title", ""), 9.5, False, "5A5A5A", editable=False)
    _line(s, "broker.subtitle", bl["broker.subtitle"], broker.get("subtitle", ""), 9.5, False, "5A5A5A", editable=False)
    _line(s, "broker.agency", bl["broker.agency"], broker.get("agency", ""), 9.5, False, "5A5A5A", editable=False)
    if broker.get("phone"):
        _line(s, "broker.phone", bl["broker.phone"], "T : %s" % broker["phone"], 9.5, True, "1A1A1A", editable=False)
    _rect(s, *_pos("bottom_bar", *L["bottom_bar"]), fill=th["bar"], name="STDp::bottom_bar")  # barre


# ───────────────────────────── Page 2 ─────────────────────────────
def _room_table(s, rooms, th, y, hhp=24.0, rhp=24.0, start_index=0, tx0=20.07, tw0=500.31, name=None):
    rows = len(rooms) + 1
    total_h = hhp + rhp * len(rooms)
    gf = s.shapes.add_table(rows, 3, Pt(tx0), Pt(y), Pt(tw0), Pt(total_h))
    _name(gf, name)
    gt = gf.table
    gt.first_row = False; gt.horz_banding = False
    gt.rows[0].height = Pt(hhp)
    for i in range(1, rows):
        gt.rows[i].height = Pt(rhp)
    gt.columns[0].width = Pt(tw0 * 0.46)
    gt.columns[1].width = Pt(tw0 * 0.27)
    gt.columns[2].width = Pt(tw0 * 0.27)

    def fill_cell(cell, text, size, bold, color, bg):
        cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(bg)
        cell.margin_top = cell.margin_bottom = Pt(1); cell.margin_left = Pt(8); cell.margin_right = Pt(4)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT; r.font.color.rgb = _rgb(color)

    for ci, htext in enumerate(["Pièce", "Étage", "Dimension"]):
        fill_cell(gt.cell(0, ci), htext, 11, True, th["th_fg"], th["th_bg"])
    for ri, room in enumerate(rooms):
        bg = th["row_alt"] if (start_index + ri) % 2 else th["row"]
        for ci in range(3):
            v = str(room[ci]) if ci < len(room) and room[ci] is not None else ""
            fill_cell(gt.cell(ri + 1, ci), v, 10, False, th["row_fg"], bg)


def _footer(s, d, th):
    L = LAYOUT
    broker = d.get("broker", {}); img = d.get("images", {}) or {}
    h2x, h2y, h2w, h2h = _pos("hero2", *L["hero2"])
    _pic_fit(s, img.get("brand_hero") or th.get("hero"), h2x, h2y, h=h2h, name="STDp::hero2")  # héros
    bl = broker2_lines(L)  # 5 lignes (texte issu du profil → position seule)
    _line(s, "broker2.name", bl["broker2.name"], broker.get("name", ""), 13, True, "1A1A1A", editable=False)
    tl = " ".join([x for x in [broker.get("title"), broker.get("subtitle")] if x])
    _line(s, "broker2.titles", bl["broker2.titles"], tl, 9, True, "1A1A1A", editable=False)
    agency = broker.get("agency", "")
    if broker.get("company"):
        agency = (agency + " | " + broker["company"]) if agency else broker["company"]
    _line(s, "broker2.agency", bl["broker2.agency"], agency, 9, False, "5A5A5A", editable=False)
    contact = [x for x in [("T : %s" % broker["phone"]) if broker.get("phone") else None,
                           ("E : %s" % broker["email"]) if broker.get("email") else None] if x]
    _line(s, "broker2.contact", bl["broker2.contact"], "  |  ".join(contact), 9, False, "5A5A5A", editable=False)
    web = broker.get("web") or broker.get("website")
    if web:
        _line(s, "broker2.web", bl["broker2.web"], "W : %s" % web, 9, False, "5A5A5A", editable=False)
    qbuf = _qr_stream(_broker_url(d, broker), th["qr"])  # QR
    if qbuf:
        qx, qy, qw, qh = _pos("qr", *L["qr"])
        _name(s.shapes.add_picture(qbuf, Pt(qx), Pt(qy), Pt(qw), Pt(qw)), "STDp::qr")
        ref = d.get("brochure_ref") or (("MLS %s" % d["mls"]) if d.get("mls") else "")
        if ref:
            _line(s, "brochure_ref", [qx - 10, qy + qw + 1, qw + 20, 12], ref, 7, False, "5A5A5A",
                  align=PP_ALIGN.CENTER, editable=False)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide2(prs, d, th):
    """Page 2 (+ page(s) de suite si trop de pièces) — même pagination que le moteur PDF."""
    L = LAYOUT
    s = _blank(prs)
    rooms = d.get("rooms", [])
    _rect(s, *_pos("headline", *L["p2_title"]), fill=th["p2_banner"],
          text=d.get("headline", d.get("title", "")), size=16, bold=True, color=th["p2_title_fg"],
          align=PP_ALIGN.LEFT, padx=12, name="STD::headline")
    if d.get("description"):
        _rect(s, *_pos("description", *L["desc"]), fill=th["desc_bg"], text=d["description"], size=12,
              color="1A1A1A", align=PP_ALIGN.JUSTIFY, anchor=MSO_ANCHOR.TOP, padx=12, name="STD::description")
    interior = (d.get("interior", []) + [None, None, None])
    for idx2 in range(3):
        _pic(s, interior[idx2], *_pos("photo.%d" % idx2, *L["photos"][idx2]), name="STDp::photo.%d" % idx2)

    if not rooms:
        _footer(s, d, th); return

    tbx, tby, tbw, tbh = _pos("table", *L["table"])
    TBL_TOP, TBL_BOT, HHP, PAGE_BOT = tby, tby + tbh, 24.0, 700.0
    n = len(rooms); avail = (TBL_BOT - TBL_TOP) - HHP
    if n * 15 <= avail:                       # tableau + pied sur la page 2 (rangées comprimées)
        rhp = min(26.0, avail / n)
        _room_table(s, rooms, th, TBL_TOP, HHP, rhp, 0, tbx, tbw, name="STD::table")
        _footer(s, d, th); return
    # Débordement : page 2 remplie sans pied ; suite + pied en page(s) suivante(s).
    rhp = 24.0
    cap2 = max(1, int((PAGE_BOT - (TBL_TOP + HHP)) / rhp))
    _room_table(s, rooms[:cap2], th, TBL_TOP, HHP, rhp, 0, tbx, tbw); idx = cap2
    footer_done = False
    while idx < n:
        s2 = _blank(prs); top = 40.0; remaining = n - idx
        if remaining <= int((TBL_BOT - (top + HHP)) / rhp):
            _room_table(s2, rooms[idx:], th, top, HHP, rhp, idx, tbx, tbw); idx = n
            _footer(s2, d, th); footer_done = True
        else:
            k = max(1, int((PAGE_BOT - (top + HHP)) / rhp))
            _room_table(s2, rooms[idx:idx + k], th, top, HHP, rhp, idx, tbx, tbw); idx += k
    if not footer_done:
        _footer(_blank(prs), d, th)


def render(data, out):
    global LAYOUT
    tpl = data.get("template") or "unifamilial"
    LAYOUT = load_layout(tpl, override=data.get("layout"))  # surcharge propriété éventuelle
    th = THEMES.get(tpl, THEMES["unifamilial"])
    prs = Presentation()
    prs.slide_width = Pt(540); prs.slide_height = Pt(720)
    slide1(prs, data, th)
    slide2(prs, data, th)
    prs.save(out)
    return out


def main():
    try:
        raw = sys.stdin.buffer.read().decode("utf-8") or "{}"
        payload = json.loads(raw)
        data = payload.get("data") or {}
        out = payload.get("out")
        if not out:
            print(json.dumps({"error": "out requis"})); return
        render(data, out)
        print(json.dumps({"path": out}))
    except Exception as e:  # noqa: BLE001
        import traceback
        print(json.dumps({"error": str(e), "trace": traceback.format_exc()[-600:]}))


if __name__ == "__main__":
    main()
