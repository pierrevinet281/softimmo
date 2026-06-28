# -*- coding: utf-8 -*-
"""Primitives PPTX natives qui MIROITENT le moteur PDF ReportLab de la brochure RPA
(8,5x11 portrait, 612x792 pt).

Modèle de coordonnées : l'appelant utilise les MÊMES coordonnées que le script PDF →
(x, y, w, h) avec y mesuré depuis le BAS de la page, exactement comme ReportLab. Ces
primitives convertissent vers l'espace EMU (haut-gauche) de PowerPoint. Porté du gabarit
ancêtre « rpa_mlt » (pptx_helpers.py) qui reproduisait le PDF au point près.

Les icônes Font Awesome sont rendues en PNG net (PowerPoint déroute les glyphes PUA des
polices d'icônes vers une police emoji couleur). Déterministe, sans IA.
"""
import os
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image, ImageDraw, ImageFont

PW, PH = 612.0, 792.0
EMU_PT = 12700

# ---- calibration (réglée contre le rendu PowerPoint) ----
ASC = {"osw": 0.78, "sg": 0.74, "sgb": 0.74}   # facteur d'ascendante (ligne de base) par famille
PARA_DY = 2.0                                   # léger décalage du haut de paragraphe (pt, vers le bas)


def E(pt):
    return Emu(int(round(pt * EMU_PT)))


def RGB(hexstr):
    if isinstance(hexstr, RGBColor):
        return hexstr
    return RGBColor.from_string(str(hexstr).lstrip("#"))


# code de police (même que le PDF) -> (typeface pptx, gras, italique, groupe)
FONTS = {
    "Osw-L":  ("Oswald", False, False, "osw"),
    "Osw-M":  ("Oswald", False, False, "osw"),
    "Osw-SB": ("Oswald", True,  False, "osw"),
    "Osw-B":  ("Oswald", True,  False, "osw"),
    "Sg":     ("Segoe UI",            False, False, "sg"),
    "Sg-L":   ("Segoe UI Light",      False, False, "sg"),
    "Sg-SL":  ("Segoe UI Semilight",  False, False, "sg"),
    "Sg-SB":  ("Segoe UI Semibold",   False, False, "sgb"),
    "Sg-B":   ("Segoe UI",            True,  False, "sgb"),
    "Sg-I":   ("Segoe UI",            False, True,  "sg"),
    "FA":     ("FA Solid",            False, False, "fa"),
    "FAB":    ("FA Brands",           False, False, "fa"),
}


# ------------------------------------------------------------------ bas niveau
def _no_autofit(tf):
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)


def _set_run(run, code, size, color):
    face, bold, ital, grp = FONTS[code]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = ital
    run.font.color.rgb = RGB(color)
    run.font.name = face
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs", "a:sym"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", face)
    return rPr


def _tracking(rPr, tracking):
    if tracking:
        rPr.set("spc", str(int(tracking * 100)))


def _name(shp, name):
    if name:
        shp.name = name
    return shp


# ------------------------------------------------------------------ texte une ligne
def text_line(slide, x, y_base, s, code, size, color, align="l", tracking=0.0, name=None):
    """Dessine une ligne de texte dont la ligne de base est à PDF (x, y_base)."""
    face, bold, ital, grp = FONTS[code]
    asc = ASC.get(grp, 0.75)
    top = (PH - y_base) - asc * size
    W = 460.0
    if align == "l":
        left = x; al = PP_ALIGN.LEFT
    elif align == "c":
        left = x - W / 2; al = PP_ALIGN.CENTER
    else:
        left = x - W; al = PP_ALIGN.RIGHT
    tb = slide.shapes.add_textbox(E(left), E(top), E(W), E(size * 1.7))
    tf = tb.text_frame; _no_autofit(tf)
    p = tf.paragraphs[0]; p.alignment = al
    try:
        p.line_spacing = 1.0
    except Exception:  # noqa: BLE001
        pass
    run = p.add_run(); run.text = "" if s is None else str(s)
    rPr = _set_run(run, code, size, color); _tracking(rPr, tracking)
    return _name(tb, name)


# ------------------------------------------------------------------ icônes en PNG net
_ICON_FONTS = {}
_SS = 5
_ASSET_DIR = None


def set_asset_dir(d):
    global _ASSET_DIR
    _ASSET_DIR = d
    os.makedirs(d, exist_ok=True)


def set_icon_fonts(solid_path, brands_path):
    if solid_path and os.path.exists(solid_path):
        _ICON_FONTS["FA"] = solid_path
    if brands_path and os.path.exists(brands_path):
        _ICON_FONTS["FAB"] = brands_path


def _color_rgb(color):
    s = str(color).lstrip("#")
    return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _icon_png(glyph_char, font, size_pt, rgb):
    cp = ord(glyph_char)
    fn = os.path.join(_ASSET_DIR, "icon_%s_%04X_%d_%02X%02X%02X.png" % (
        font, cp, int(round(size_pt * 10)), rgb[0], rgb[1], rgb[2]))
    if os.path.exists(fn):
        with Image.open(fn) as im:
            return fn, im.width / _SS, im.height / _SS
    px = max(8, int(round(size_pt * _SS)))
    f = ImageFont.truetype(_ICON_FONTS[font], px)
    d0 = ImageDraw.Draw(Image.new("RGBA", (px * 3, px * 3), (0, 0, 0, 0)))
    bb = d0.textbbox((0, 0), glyph_char, font=f)
    w = max(1, bb[2] - bb[0]); h = max(1, bb[3] - bb[1]); pad = 2
    img = Image.new("RGBA", (w + 2 * pad, h + 2 * pad), (0, 0, 0, 0))
    ImageDraw.Draw(img).text((pad - bb[0], pad - bb[1]), glyph_char, font=f, fill=(rgb[0], rgb[1], rgb[2], 255))
    img.save(fn)
    return fn, img.width / _SS, img.height / _SS


def icon(slide, glyph_char, cx, cy, size, color, font="FA"):
    if font not in _ICON_FONTS or not glyph_char:
        return None
    rgb = _color_rgb(color)
    path, wpt, hpt = _icon_png(glyph_char, font, size, rgb)
    return picture_raw(slide, path, cx - wpt / 2, cy - hpt / 2, wpt, hpt)


# ------------------------------------------------------------------ paragraphes (avec retour à la ligne)
def para(slide, x, y_top, w, text, code, size, color, leading, align="l", name=None):
    top = (PH - y_top) + PARA_DY
    al = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT, "j": PP_ALIGN.JUSTIFY}[align]
    tb = slide.shapes.add_textbox(E(x), E(top), E(w), E(200))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]; p.alignment = al
    p.line_spacing = Pt(leading)
    run = p.add_run(); run.text = "" if text is None else str(text)
    _set_run(run, code, size, color)
    return _name(tb, name)


# ------------------------------------------------------------------ formes
def _solid(shp, color):
    shp.fill.solid(); shp.fill.fore_color.rgb = RGB(color)


def _alpha_on_fill(shp, alpha):
    sp = shp.fill._xPr.find(qn("a:solidFill"))
    clr = sp.find(qn("a:srgbClr"))
    a = clr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))}); clr.append(a)


def _no_line(shp):
    shp.line.fill.background()


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.0, fill_alpha=None, name=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0 else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, E(x), E(PH - y - h), E(w), E(h))
    shp.shadow.inherit = False
    if radius > 0:
        try:
            shp.adjustments[0] = max(0.0, min(0.5, radius / min(w, h)))
        except Exception:  # noqa: BLE001
            pass
    if fill is None:
        shp.fill.background()
    else:
        _solid(shp, fill)
        if fill_alpha is not None:
            _alpha_on_fill(shp, fill_alpha)
    if line is None:
        _no_line(shp)
    else:
        shp.line.color.rgb = RGB(line); shp.line.width = Pt(line_w)
    return _name(shp, name)


def oval(slide, cx, cy, r, fill=None, line=None, line_w=1.0, fill_alpha=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, E(cx - r), E((PH - cy) - r), E(2 * r), E(2 * r))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        _solid(shp, fill)
        if fill_alpha is not None:
            _alpha_on_fill(shp, fill_alpha)
    if line is None:
        _no_line(shp)
    else:
        shp.line.color.rgb = RGB(line); shp.line.width = Pt(line_w)
    return shp


def line(slide, x1, y1, x2, y2, color, w=1.0):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(x1), E(PH - y1), E(x2), E(PH - y2))
    cn.line.color.rgb = RGB(color); cn.line.width = Pt(w)
    cn.shadow.inherit = False
    return cn


# ------------------------------------------------------------------ dégradés
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _gradfill_xml(stops, angle_deg):
    ang = int((angle_deg % 360) * 60000)
    gs = []
    for pos, rgb, alpha in stops:
        r, g, b = rgb
        a = "" if alpha is None else '<a:alpha val="%d"/>' % int(alpha * 100000)
        gs.append('<a:gs pos="%d"><a:srgbClr val="%02X%02X%02X">%s</a:srgbClr></a:gs>' % (int(pos * 100000), r, g, b, a))
    xml = ('<a:gradFill xmlns:a="%s" rotWithShape="1"><a:gsLst>' % A_NS
           + "".join(gs)
           + '</a:gsLst><a:lin ang="%d" scaled="1"/></a:gradFill>' % ang)
    return etree.fromstring(xml)


def _apply_grad(shp, stops, angle_deg):
    spPr = shp._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        e = spPr.find(qn(tag))
        if e is not None:
            spPr.remove(e)
    grad = _gradfill_xml(stops, angle_deg)
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)


def _rgbtuple(c):
    s = str(c).lstrip("#"); return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def grad_rect(slide, x, y, w, h, c1, c2, radius=0.0, vertical=True):
    shp = rect(slide, x, y, w, h, fill="#FFFFFF", line=None, radius=radius)
    a = c1 if isinstance(c1, tuple) else _rgbtuple(c1)
    b = c2 if isinstance(c2, tuple) else _rgbtuple(c2)
    angle = 90 if vertical else 0
    _apply_grad(shp, [(0.0, a, None), (1.0, b, None)], angle)
    return shp


def scrim(slide, x, y, w, h, top_alpha=0, bot_alpha=210, base=(8, 30, 40)):
    shp = rect(slide, x, y, w, h, fill="#000000", line=None, radius=0)
    mid = 0.5
    am = top_alpha + (bot_alpha - top_alpha) * (mid ** 1.6)
    stops = [(0.0, base, top_alpha / 255.0),
             (mid, base, am / 255.0),
             (1.0, base, bot_alpha / 255.0)]
    _apply_grad(shp, stops, 90)
    return shp


# ------------------------------------------------------------------ images (cover-crop + arrondi)
def _cover_crop(img, tw, th):
    w, h = img.size; ta = tw / th; a = w / h
    if a > ta:
        nw = int(h * ta); x = (w - nw) // 2; box = (x, 0, x + nw, h)
    else:
        nh = int(w / ta); y = (h - nh) // 2; box = (0, y, w, y + nh)
    return img.crop(box).resize((tw, th), Image.LANCZOS)


_piccache = {}


def _cropped_asset(path, w_pt, h_pt, dpi=240):
    key = (path, round(w_pt, 1), round(h_pt, 1), dpi)
    if key in _piccache:
        return _piccache[key]
    tw = max(2, int(w_pt / 72 * dpi)); th = max(2, int(h_pt / 72 * dpi))
    img = Image.open(path).convert("RGB")
    tw = min(tw, int(img.size[0] * 1.05)); th = min(th, int(img.size[1] * 1.05))
    out = _cover_crop(img, tw, th)
    base = os.path.splitext(os.path.basename(path))[0]
    fn = os.path.join(_ASSET_DIR, "%s_%dx%d.jpg" % (base, tw, th))
    out.save(fn, "JPEG", quality=88, optimize=True)
    _piccache[key] = fn
    return fn


def _round_pic(pic, radius_pt, w_pt, h_pt):
    spPr = pic._element.spPr
    for tag in ("a:prstGeom", "a:custGeom"):
        geo = spPr.find(qn(tag))
        if geo is not None:
            spPr.remove(geo)
    adj = max(0.0, min(0.5, radius_pt / min(w_pt, h_pt)))
    xml = ('<a:prstGeom xmlns:a="%s" prst="roundRect"><a:avLst>'
           '<a:gd name="adj" fmla="val %d"/></a:avLst></a:prstGeom>' % (A_NS, int(adj * 100000)))
    el = etree.fromstring(xml)
    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is not None:
        xfrm.addnext(el)
    else:
        spPr.insert(0, el)


def picture(slide, path, x, y, w, h, radius=10.0, cover=True, dpi=240):
    src = _cropped_asset(path, w, h, dpi) if cover else path
    pic = slide.shapes.add_picture(src, E(x), E(PH - y - h), E(w), E(h))
    pic.shadow.inherit = False
    if radius > 0:
        _round_pic(pic, radius, w, h)
    return pic


def picture_raw(slide, path, x, y, w, h):
    pic = slide.shapes.add_picture(path, E(x), E(PH - y - h), E(w), E(h))
    pic.shadow.inherit = False
    return pic


def img_size(path):
    return Image.open(path).size
