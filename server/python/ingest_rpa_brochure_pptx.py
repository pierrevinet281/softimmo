# -*- coding: utf-8 -*-
"""Ingestion d'un PPTX RPA édité → contenu RPA (aller-retour, docs/09).

Relit les formes NOMMÉES « RPA::<chemin.pointé> » écrites par render_rpa_brochure_pptx.py et
superpose le texte édité sur le CONTENU DE BASE fourni (défauts + surcharge live). Ainsi, seules
les chaînes éditables sont mises à jour : les ICÔNES et la structure (longueurs de tableaux,
clés non textuelles) sont PRÉSERVÉES. Déterministe, sans IA.

E/S : lit {"pptx": "<chemin>", "base": {<contenu d'origine>}} sur stdin ;
renvoie {"content": {<contenu mis à jour>}}.
"""
import sys
import json
import copy

PH = 792.0       # hauteur lettre (pt)
EMU_PT = 12700   # 1 pt = 12700 EMU


def _text(shape):
    """Texte d'une forme (paragraphes non vides joints par des sauts de ligne)."""
    try:
        if not shape.has_text_frame:
            return None
    except Exception:  # noqa: BLE001
        return None
    lines = []
    for p in shape.text_frame.paragraphs:
        txt = "".join(r.text for r in p.runs)
        lines.append(txt)
    # retire les lignes vides en tête/queue, conserve les sauts internes
    while lines and not lines[0].strip():
        lines.pop(0)
    while lines and not lines[-1].strip():
        lines.pop()
    return "\n".join(lines)


def _set_path(root, dotted, value):
    """Affecte value au chemin pointé dans root (existant). Les segments numériques = index."""
    parts = dotted.split(".")
    cur = root
    for i, seg in enumerate(parts):
        last = i == len(parts) - 1
        idx = int(seg) if seg.isdigit() else None
        if last:
            if idx is not None:
                if isinstance(cur, list) and 0 <= idx < len(cur):
                    cur[idx] = value
            elif isinstance(cur, dict):
                cur[seg] = value
            return
        # descente
        if idx is not None:
            if isinstance(cur, list) and 0 <= idx < len(cur):
                cur = cur[idx]
            else:
                return  # structure inattendue → on ignore (garde-fou)
        else:
            if isinstance(cur, dict) and isinstance(cur.get(seg), (dict, list)):
                cur = cur[seg]
            else:
                return


def _box_pt(sh):
    """Boîte de la forme en pt, origine BAS-gauche (comme ReportLab) : [x, y_bas, w, h]."""
    try:
        x = float(sh.left) / EMU_PT; top = float(sh.top) / EMU_PT
        w = float(sh.width) / EMU_PT; h = float(sh.height) / EMU_PT
    except Exception:  # noqa: BLE001
        return None
    return [round(x, 2), round(PH - top - h, 2), round(w, 2), round(h, 2)]


def build_content(prs, base):
    content = copy.deepcopy(base) if isinstance(base, dict) else {}
    layout = {}
    for slide in prs.slides:
        for sh in slide.shapes:
            name = getattr(sh, "name", "") or ""
            if not name.startswith("RPA::") or name.startswith("RPA::__"):
                continue
            path = name[len("RPA::"):]
            if not path:
                continue
            # Texte édité → superposé sur le contenu (les formes image n'ont pas de texte).
            val = _text(sh)
            if val is not None and val != "":
                _set_path(content, path, val)
            # Position de la forme → override de layout (aller-retour des positions, Phase C).
            box = _box_pt(sh)
            if box:
                layout[path] = box
    return content, layout


def _emit(obj):
    sys.stdout.buffer.write(json.dumps(obj, ensure_ascii=False).encode("utf-8"))
    sys.stdout.buffer.flush()


def main():
    try:
        from pptx import Presentation
        payload = json.loads(sys.stdin.buffer.read().decode("utf-8") or "{}")
        pptx = payload.get("pptx")
        if not pptx:
            _emit({"error": "pptx requis"}); return
        prs = Presentation(pptx)
        content, layout = build_content(prs, payload.get("base") or {})
        _emit({"content": content, "layout": layout})
    except Exception as e:  # noqa: BLE001
        import traceback
        _emit({"error": str(e), "trace": traceback.format_exc()[-700:]})


if __name__ == "__main__":
    main()
