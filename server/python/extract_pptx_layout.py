# -*- coding: utf-8 -*-
"""Extraction des positions d'un gabarit PowerPoint (.pptx) → JSON normalisé.

Sert de socle au « round-trip » des brochures (docs/09) : le courtier ajuste les
positions dans PowerPoint, on ré-extrait ici, puis on reporte les coordonnées dans
`render_brochure.py` (espace de référence 540×720 pt = 7,5×10 po).

Résout les transformations de groupe pour donner des coordonnées ABSOLUES (slide).
Unités : points PostScript (1 po = 72 pt ; 914400 EMU = 1 po).

Usage :
    python extract_pptx_layout.py "chemin/brochure.pptx" [> layout.json]
"""
import sys
import json

from pptx import Presentation

EMU = 914400.0


def pt(v):
    return None if v is None else round(v / EMU * 72, 2)


def _xfrm(el):
    """Retourne (off_x, off_y, ext_cx, ext_cy, choff_x, choff_y, chext_cx, chext_cy) en EMU."""
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    x = el.find(".//%sxfrm" % ns)
    if x is None:
        return None
    off = x.find("%soff" % ns); ext = x.find("%sext" % ns)
    choff = x.find("%schOff" % ns); chext = x.find("%schExt" % ns)
    g = lambda n, a: int(n.get(a)) if n is not None else None
    return (g(off, "x"), g(off, "y"), g(ext, "cx"), g(ext, "cy"),
            g(choff, "x"), g(choff, "y"), g(chext, "cx"), g(chext, "cy"))


def _abs_pos(sh, tf):
    """Position absolue (slide) d'une shape, en appliquant la transform de groupe `tf`."""
    l, t, w, h = sh.left, sh.top, sh.width, sh.height
    if tf is None or None in (l, t):
        return l, t, w, h
    ox, oy, ecx, ecy, cox, coy, cecx, cecy = tf
    sx = ecx / cecx if cecx else 1
    sy = ecy / cecy if cecy else 1
    al = ox + (l - (cox or 0)) * sx
    at = oy + (t - (coy or 0)) * sy
    return al, at, w * sx, h * sy


def _hex(color):
    try:
        return "#" + str(color.rgb)
    except Exception:
        return None


def _shape(sh, tf=None):
    l, t, w, h = _abs_pos(sh, tf)
    out = {"name": sh.name, "type": int(sh.shape_type) if sh.shape_type else None,
           "x": pt(l), "y": pt(t), "w": pt(w), "h": pt(h)}
    try:
        if sh.fill.type == 1:
            out["fill"] = _hex(sh.fill.fore_color)
    except Exception:
        pass
    try:
        if sh.has_text_frame and sh.text_frame.text:
            out["text"] = sh.text_frame.text
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    out["font"] = {"size": pt(r.font.size), "bold": r.font.bold,
                                   "color": _hex(r.font.color) if r.font.color and r.font.color.type else None,
                                   "name": r.font.name}
                    break
                if "font" in out:
                    break
    except Exception:
        pass
    if int(sh.shape_type or 0) == 6:  # GROUP
        ctf = _xfrm(sh._element)
        out["children"] = [_shape(ch, ctf) for ch in sh.shapes]
    return out


def extract(path):
    prs = Presentation(path)
    return {"slide_w": pt(prs.slide_width), "slide_h": pt(prs.slide_height),
            "slides": [[_shape(sh) for sh in s.shapes] for s in prs.slides]}


if __name__ == "__main__":
    print(json.dumps(extract(sys.argv[1]), ensure_ascii=False, indent=2))
