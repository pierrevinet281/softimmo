# -*- coding: utf-8 -*-
"""Jumeau PPTX éditable de l'offre de services (Module 3 — aller-retour, docs/09).

Le PDF de l'offre est un document à FLUX (Platypus) ; pour permettre l'édition dans
PowerPoint puis la resynchronisation, on génère ici une diapo ÉDITABLE par section
(titre + corps à puces), avec des formes NOMMÉES (« OFF::<clé>::<type>::<partie> ») que
`ingest_offre_pptx.py` relit pour reconstruire le contenu de l'offre. Déterministe, sans IA.

E/S : lit {"data": {...}, "out": "<chemin.pptx>"} sur stdin, écrit le PPTX, renvoie {"path"}.
Conformité : mentions agence + courtier (désignation) en pied de chaque diapo (LCI/OACIQ).
"""
import os
import sys
import json

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

try:
    from PIL import Image as PILImage
except Exception:  # noqa: BLE001
    PILImage = None

FONT = "Segoe UI"
FONT_BOLD = "Segoe UI Semibold"
PW, PH = 612.0, 792.0   # Lettre portrait (pt), cohérent avec le PDF
M = 50.0
CW = PW - 2 * M

INK = "1A1A1A"
INK2 = "5A5A5A"
INK3 = "8A8A8A"
WHITE = "FFFFFF"
LINE = "D7DEEE"


def _rgb(h):
    return RGBColor.from_string(h.lstrip("#"))


def _lum(h):
    h = h.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return (0.2126 * r + 0.7152 * g + 0.0722 * b) / 255.0


def _on(h):
    return INK if _lum(h) > 0.62 else WHITE


def _box(slide, x, y, w, h, name=None):
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    if name:
        tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def _para(tf, text, *, level=0, size=11, bold=False, italic=False, color=INK, first=False, align=None):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.level = max(0, min(4, level))
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = "" if text is None else str(text)
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.italic = bool(italic)
    r.font.name = FONT_BOLD if bold else FONT
    r.font.color.rgb = _rgb(color)
    return p


# ── Type structurel d'une section (encodé dans le nom de forme, lu à l'ingestion) ──
def _stype(key, sec):
    if (sec.get("kind") == "asset") or sec.get("image"):
        return "asset"
    if key in ("why", "services"):
        return "groups"
    if key == "marketing":
        return "subgroups"
    if key in ("opportunities", "value_add"):
        return "list"
    if key == "timeline":
        return "steps"
    if key == "testimonials":
        return "testi"
    if key in ("guarantee", "fees", "next_steps"):
        return "body"
    if sec.get("groups"):
        return "groups"
    if sec.get("items"):
        return "list"
    return "body"


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _footer(slide, broker):
    desig = " ".join([s for s in [broker.get("title"), broker.get("subtitle")] if s])
    agency = broker.get("agency", "")
    if broker.get("company"):
        agency = (agency + " | " + broker["company"]) if agency else broker["company"]
    bits = [broker.get("name"), desig, agency]
    tf = _box(slide, M, PH - 30, CW, 20)
    _para(tf, "  ·  ".join([b for b in bits if b]), size=7.5, color=INK2, first=True)
    _para(tf, "Document de présentation — ne constitue pas un contrat de courtage.", size=6.6, color=INK3)


def _brand_slide(prs, data, content):
    broker = data.get("broker", {})
    th = data.get("theme") or {}
    band = (th.get("band_color") or "#314897").lstrip("#")
    fg = _on(band)
    slide = _blank(prs)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(0), Pt(0), Pt(PW), Pt(96))
    rect.fill.solid(); rect.fill.fore_color.rgb = _rgb(band); rect.line.fill.background()
    rect.name = "OFF::__band"
    logo = data.get("logo")
    tx = M
    if logo and os.path.exists(logo) and PILImage is not None:
        try:
            iw, ih = PILImage.open(logo).size
            h = 40.0; w = h * iw / ih
            slide.shapes.add_picture(logo, Pt(M), Pt(28), height=Pt(h))
            tx = M + w + 16
        except Exception:  # noqa: BLE001
            pass
    tf = _box(slide, tx, 18, PW - tx - M, 70, name="OFF::__title")
    _para(tf, content.get("doc_title", "Proposition d'offre de services"), size=18, bold=True, color=fg, first=True)
    if content.get("subtitle"):
        _para(tf, content["subtitle"], size=11, color=fg)
    desig = " ".join([s for s in [broker.get("title"), broker.get("subtitle")] if s])
    _para(tf, "%s — %s" % (broker.get("name", ""), desig), size=10, bold=True, color=fg)
    # méta
    mf = _box(slide, M, 112, CW, 60)
    first = True
    if data.get("client_name"):
        _para(mf, "Préparée pour : %s" % data["client_name"], size=10.5, color=INK2, first=first); first = False
    _para(mf, "Préparée par : %s" % broker.get("name", ""), size=10.5, color=INK2, first=first)
    if data.get("date"):
        _para(mf, "Date : %s" % data["date"], size=10.5, color=INK2)
    _footer(slide, broker)


def _section_slide(prs, data, key, sec):
    broker = data.get("broker", {})
    th = data.get("theme") or {}
    title_color = (th.get("title_color") or "#314897").lstrip("#")
    stype = _stype(key, sec)
    slide = _blank(prs)

    # Section image (asset) : photo + légende.
    if stype == "asset":
        img = sec.get("image")
        if img and os.path.exists(img):
            w, h = (CW, CW * 0.6)
            if PILImage is not None:
                try:
                    iw, ih = PILImage.open(img).size
                    h = min(360.0, CW * ih / iw); w = h * iw / ih
                    if w > CW:
                        w = CW; h = w * ih / iw
                except Exception:  # noqa: BLE001
                    pass
            pic = slide.shapes.add_picture(img, Pt(M + (CW - w) / 2), Pt(70), width=Pt(w), height=Pt(h))
            pic.name = "OFF::%s::asset::img" % key
        cf = _box(slide, M, 70 + 366, CW, 24, name="OFF::%s::asset::cap" % key)
        _para(cf, sec.get("caption", ""), size=9.5, italic=True, color=INK2, first=True)
        _footer(slide, broker)
        return

    # Titre de section.
    hf = _box(slide, M, 44, CW, 34, name="OFF::%s::%s::head" % (key, stype))
    _para(hf, sec.get("heading", ""), size=18, bold=True, color=title_color, first=True)

    # Corps (puces, niveaux pour les groupes).
    bf = _box(slide, M, 86, CW, PH - 86 - 40, name="OFF::%s::%s::body" % (key, stype))
    bf.vertical_anchor = MSO_ANCHOR.TOP
    first = True

    def add(text, **kw):
        nonlocal first
        _para(bf, text, first=first, **kw); first = False

    if sec.get("intro"):
        add(sec["intro"], size=11, color=INK)
    if stype in ("groups", "subgroups"):
        groups = sec.get("groups") or sec.get("subgroups") or []
        for g in groups:
            add(g.get("label", ""), level=0, size=12, bold=True, color=title_color)
            for it in (g.get("items") or []):
                add(it, level=1, size=10.5, color=INK)
    elif stype == "list":
        for it in (sec.get("items") or []):
            add(it, level=0, size=10.5, color=INK)
    elif stype == "steps":
        for s in (sec.get("steps") or []):
            add("%s — %s" % (s.get("label", ""), s.get("text", "")), level=0, size=10.5, color=INK)
    elif stype == "testi":
        for tt in (sec.get("items") or []):
            add("« %s » — %s" % (tt.get("quote", ""), tt.get("author", "")), level=0, size=10.5, color=INK)
    else:  # body
        if sec.get("body"):
            add(sec["body"], size=11, color=INK)
        if sec.get("note"):
            add(sec["note"], size=9, italic=True, color=INK2)
    if first:  # corps vide → garder une ligne éditable
        _para(bf, "", first=True)
    _footer(slide, broker)


def render(data, out):
    content = (data.get("content") or {})
    # langue unique pour le PPTX : la première fournie (bi → fr).
    langs = data.get("langs") or ["fr"]
    lang = langs[0] if langs[0] in content else (next(iter(content), "fr"))
    c = content.get(lang) or {}

    prs = Presentation()
    prs.slide_width = Pt(PW)
    prs.slide_height = Pt(PH)

    _brand_slide(prs, data, c)
    order = c.get("sections")
    if not isinstance(order, list) or not order:
        order = [{"key": k} for k in ["why", "guarantee", "services", "marketing", "opportunities",
                                      "timeline", "fees", "value_add", "testimonials", "next_steps"] if c.get(k)]
    for s in order:
        if not isinstance(s, dict) or s.get("hidden"):
            continue
        key = s.get("key")
        if key and (c.get(key) or s.get("custom")):
            _section_slide(prs, data, key, c.get(key) or {})

    prs.save(out)
    return out


def main():
    try:
        raw = sys.stdin.buffer.read().decode("utf-8") or "{}"
        payload = json.loads(raw)
        data = payload.get("data") or {}
        out = payload.get("out")
        if not out:
            print(json.dumps({"error": "out requis"})); return
        render(data, out)
        print(json.dumps({"path": out}))
    except Exception as e:  # noqa: BLE001
        import traceback
        print(json.dumps({"error": str(e), "trace": traceback.format_exc()[-700:]}))


if __name__ == "__main__":
    main()
